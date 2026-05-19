"""
ATT-branded PPT: 横版2D闯关小游戏项目汇报
Brand: #2E74B6 (blue), #D4EDFE (light blue), #FFD116 (yellow), #FFFFFF
Fonts: Tahoma Bold (titles), Calibri (body)
Size: 20 x 11.25 inches (16:9)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colors ──────────────────────────────────────────────────────
ATT_BLUE   = RGBColor(0x2E, 0x74, 0xB6)
ATT_LBLUE  = RGBColor(0xD4, 0xED, 0xFE)
ATT_YELLOW = RGBColor(0xFF, 0xD1, 0x16)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BLUE  = RGBColor(0x1A, 0x47, 0x7A)
GRAY_TEXT  = RGBColor(0x44, 0x44, 0x44)
LIGHT_GRAY = RGBColor(0xF4, 0xF7, 0xFB)

W = Inches(20)
H = Inches(11.25)
FOOTER_TXT = "Proprietary and confidential. Please do not circulate without consent from ATT Group"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # blank

# ── Helper functions ─────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, x, y, w, h, text, font_name, font_size, bold, color,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_multiline_textbox(slide, x, y, w, h, lines, font_name, font_size, bold, color,
                          align=PP_ALIGN.LEFT, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox

def add_footer(slide):
    add_textbox(slide, 0.5, 10.75, 19, 0.4,
                FOOTER_TXT, 'Calibri', 8, False, RGBColor(0x99,0x99,0x99),
                align=PP_ALIGN.CENTER)

def add_card(slide, x, y, w, h, title, bullets, accent_color):
    """Card with colored top strip, title, and bullet list."""
    # Card background
    card = add_rect(slide, x, y, w, h, LIGHT_GRAY)
    # Accent top strip
    add_rect(slide, x, y, w, 0.12, accent_color)
    # Title
    add_textbox(slide, x+0.18, y+0.18, w-0.36, 0.45,
                title, 'Tahoma', 16, True, DARK_BLUE)
    # Bullets
    add_multiline_textbox(slide, x+0.18, y+0.7, w-0.36, h-0.9,
                          bullets, 'Calibri', 13, False, GRAY_TEXT)

def add_stat_callout(slide, x, y, w, h, number, label, bg_color):
    add_rect(slide, x, y, w, h, bg_color)
    add_textbox(slide, x, y+0.15, w, 0.65,
                number, 'Tahoma', 36, True, WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, y+0.75, w, 0.4,
                label, 'Calibri', 12, False, WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1: Cover
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
# Full blue background
add_rect(slide, 0, 0, 20, 11.25, ATT_BLUE)
# Yellow accent left strip
add_rect(slide, 0, 0, 0.25, 11.25, ATT_YELLOW)
# Light blue right panel
add_rect(slide, 13.5, 0, 6.5, 11.25, DARK_BLUE)

# Game icon area (placeholder circle)
dot = slide.shapes.add_shape(9, Inches(14.5), Inches(2.5), Inches(4.5), Inches(4.5))
dot.fill.solid(); dot.fill.fore_color.rgb = ATT_YELLOW
dot.line.fill.background()
# "🎮" label inside circle
add_textbox(slide, 14.5, 3.8, 4.5, 1.5,
            '🎮', 'Calibri', 72, False, DARK_BLUE, align=PP_ALIGN.CENTER)

# Title area
add_textbox(slide, 0.9, 2.2, 12.2, 1.2,
            '平台冒险 · 肉鸽养成版',
            'Tahoma', 48, True, WHITE)
add_textbox(slide, 0.9, 3.6, 12.2, 0.6,
            'Platform Adventure · Roguelite Edition',
            'Calibri', 22, False, ATT_LBLUE)
add_textbox(slide, 0.9, 4.5, 12.2, 0.5,
            '项目汇报 · Work & Learning Review',
            'Calibri', 18, False, ATT_YELLOW)

# Divider line
add_rect(slide, 0.9, 5.3, 10, 0.04, ATT_YELLOW)

# Metadata
add_textbox(slide, 0.9, 5.6, 12, 0.4,
            'Kenny  ·  ATT Group  ·  2026',
            'Calibri', 16, False, RGBColor(0xCC,0xDD,0xEE))

# Footer
add_textbox(slide, 0.5, 10.65, 19, 0.4,
            FOOTER_TXT, 'Calibri', 8, False, RGBColor(0x88,0xAA,0xCC),
            align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2: Project Overview
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
# White bg
add_rect(slide, 0, 0, 20, 11.25, WHITE)
# Top blue header band
add_rect(slide, 0, 0, 20, 1.35, ATT_BLUE)
add_rect(slide, 0, 0, 0.25, 1.35, ATT_YELLOW)

add_textbox(slide, 0.6, 0.3, 18, 0.75,
            '项目概述', 'Tahoma', 32, True, WHITE)

# 3 overview cards
add_card(slide, 0.5, 1.6, 5.8, 4.5, '🎯 项目定位',
         ['横版2D闯关游戏', '融合 Roguelite 元成长机制',
          '纯浏览器运行，无需安装', '同步包含节奏音游模块'], ATT_BLUE)

add_card(slide, 7.1, 1.6, 5.8, 4.5, '⚙️ 技术选型',
         ['HTML5 Canvas 2D 渲染', '原生 JavaScript（零依赖）',
          'localStorage 持久化存档', 'JSON 关卡配置文件'], ATT_YELLOW)

add_card(slide, 13.7, 1.6, 5.8, 4.5, '🏆 项目目标',
         ['独立实现完整游戏循环', '探索 AI 辅助开发工作流',
          '锻炼前端工程化能力', '输出可演示的完整作品'], ATT_BLUE)

# Bottom stats bar
add_rect(slide, 0, 6.4, 20, 2.0, ATT_LBLUE)
stats = [
    ('5', '难度等级'),
    ('2', '游戏模式'),
    ('∞', '随机关卡'),
    ('6', '可解锁技能'),
    ('0', '外部依赖'),
]
for i, (num, lbl) in enumerate(stats):
    add_stat_callout(slide, 0.6 + i*3.8, 6.55, 3.2, 1.65, num, lbl, ATT_BLUE)

add_footer(slide)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3: Divider — 核心功能
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 20, 11.25, DARK_BLUE)
add_rect(slide, 0, 0, 0.25, 11.25, ATT_YELLOW)
add_textbox(slide, 1.2, 3.8, 17, 1.3,
            '核心功能', 'Tahoma', 54, True, WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 1.2, 5.4, 17, 0.6,
            'Core Game Mechanics', 'Calibri', 24, False, ATT_LBLUE, align=PP_ALIGN.CENTER)
add_footer(slide)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4: Core Mechanics (2x3 grid)
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 20, 11.25, WHITE)
add_rect(slide, 0, 0, 20, 1.35, ATT_BLUE)
add_rect(slide, 0, 0, 0.25, 1.35, ATT_YELLOW)
add_textbox(slide, 0.6, 0.3, 18, 0.75, '核心游戏机制', 'Tahoma', 32, True, WHITE)

mechanics = [
    ('🗺️ 程序化关卡生成',
     ['每局随机生成平台布局', '主路径 + 分支路径设计',
      '关卡宽度随难度线性增加', '支持 JSON 预设关卡']),
    ('⚔️ 双模式战斗系统',
     ['近战：范围攻击68px', '远程：子弹速度9px/帧',
      '攻击冷却防止无限连击', '击退 + 伤害数字特效']),
    ('🎮 5档难度系统',
     ['D1新手 → D5地狱', '敌人数量/速度/血量递增',
      'D4/D5 需解锁二段跳', '各难度独立背景配色']),
    ('📈 元成长系统',
     ['跨局经验值积累', '技能点解锁永久强化',
      '成就系统激励探索', 'localStorage 跨会话保存']),
    ('🦅 飞行敌人系统',
     ['地面巡逻 + 空中飞行双类型', '各难度飞行敌数量不同',
      '飞行敌独立移动逻辑', '难度5最多8只飞行敌']),
    ('☁️ 视觉与体验',
     ['14朵视差云背景', '粒子特效死亡动画',
      '平滑镜头跟随玩家', '难度独立配色主题']),
]

cols = [0.4, 6.9, 13.4]
rows = [1.55, 5.2]
for i, (title, bullets) in enumerate(mechanics):
    c, r = cols[i % 3], rows[i // 3]
    add_card(slide, c, r, 6.0, 3.4, title, bullets, ATT_BLUE if i % 2 == 0 else ATT_YELLOW)

add_footer(slide)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5: Work Done — 工作内容
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 20, 11.25, WHITE)
add_rect(slide, 0, 0, 20, 1.35, ATT_BLUE)
add_rect(slide, 0, 0, 0.25, 1.35, ATT_YELLOW)
add_textbox(slide, 0.6, 0.3, 18, 0.75, '工作内容', 'Tahoma', 32, True, WHITE)

# Timeline-style layout
phases = [
    ('1', '游戏框架搭建', ATT_BLUE,
     '设计Canvas渲染循环 · 实现物理引擎（重力/碰撞）\n玩家移动控制 · 游戏状态机（play/over/win）'),
    ('2', '关卡生成系统', ATT_YELLOW,
     '设计程序化平台生成算法 · 主路径脊骨 + 随机分支\nJSON关卡配置支持 · 难度参数化配置系统'),
    ('3', '战斗与敌人系统', ATT_BLUE,
     '近战/远程双模式切换 · 飞行/地面双敌人类型\n伤害计算 · 击退物理 · AI巡逻逻辑'),
    ('4', '元游戏与存档', ATT_YELLOW,
     '跨局经验值系统 · 技能树设计（6种技能）\n成就解锁逻辑 · localStorage持久化'),
    ('5', '节奏游戏模块', ATT_BLUE,
     '独立rhythm.html · MIDI音频分析\n游戏中心index.html整合两个游戏'),
]

for i, (num, title, color, desc) in enumerate(phases):
    y = 1.55 + i * 1.7
    # Number circle
    circ = slide.shapes.add_shape(9, Inches(0.4), Inches(y), Inches(0.9), Inches(0.9))
    circ.fill.solid(); circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    add_textbox(slide, 0.4, y+0.08, 0.9, 0.7, num, 'Tahoma', 22, True, WHITE, align=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, 1.5, y, 5.5, 0.5, title, 'Tahoma', 17, True, DARK_BLUE)
    # Description
    add_textbox(slide, 1.5, y+0.48, 17.5, 0.8, desc, 'Calibri', 12.5, False, GRAY_TEXT)
    # Connector line (except last)
    if i < len(phases)-1:
        add_rect(slide, 0.79, y+0.9, 0.03, 0.8, RGBColor(0xCC,0xCC,0xCC))

add_footer(slide)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6: Key Learnings — 学习收获
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 20, 11.25, WHITE)
add_rect(slide, 0, 0, 20, 1.35, ATT_BLUE)
add_rect(slide, 0, 0, 0.25, 1.35, ATT_YELLOW)
add_textbox(slide, 0.6, 0.3, 18, 0.75, '学习收获', 'Tahoma', 32, True, WHITE)

# Left: Technical learnings
add_rect(slide, 0.4, 1.55, 9.3, 8.8, LIGHT_GRAY)
add_rect(slide, 0.4, 1.55, 9.3, 0.12, ATT_BLUE)
add_textbox(slide, 0.6, 1.7, 9.0, 0.55, '💻  技术能力提升', 'Tahoma', 18, True, DARK_BLUE)

tech_learnings = [
    ('Canvas 2D 渲染', '掌握 requestAnimationFrame 游戏循环、\n精灵绘制、视口裁剪优化'),
    ('物理与碰撞', 'AABB碰撞检测实现、重力模拟、\n平台边缘检测与修正'),
    ('程序化生成', '随机种子关卡算法设计、\n参数化难度配置'),
    ('状态管理', '无框架实现完整游戏状态机、\n元数据持久化存档'),
    ('AI 辅助开发', '与 Claude Code 协作迭代功能、\n学会拆解需求与验证逻辑'),
]
for i, (title, desc) in enumerate(tech_learnings):
    y = 2.4 + i * 1.52
    add_textbox(slide, 0.65, y, 8.8, 0.38, f'▶  {title}', 'Tahoma', 14, True, ATT_BLUE)
    add_textbox(slide, 0.65, y+0.38, 8.8, 0.75, desc, 'Calibri', 12, False, GRAY_TEXT)

# Right: Soft/process learnings
add_rect(slide, 10.3, 1.55, 9.3, 8.8, LIGHT_GRAY)
add_rect(slide, 10.3, 1.55, 9.3, 0.12, ATT_YELLOW)
add_textbox(slide, 10.5, 1.7, 9.0, 0.55, '🌱  工程思维与方法论', 'Tahoma', 18, True, DARK_BLUE)

soft_learnings = [
    ('从想法到产品', '完整经历：概念→设计→实现\n→调试→迭代的产品开发闭环'),
    ('模块化设计思维', '将游戏拆分为独立模块（关卡/战斗/\n存档），降低耦合，便于维护'),
    ('调试与问题定位', '学会用 console + 可视化辅助线\n快速定位碰撞和渲染 bug'),
    ('版本迭代意识', '每次功能增量提交，小步快跑，\n及时验证而非大爆炸开发'),
    ('文档与规范意识', '配置化设计（CFG对象）让代码\n自解释，减少硬编码'),
]
for i, (title, desc) in enumerate(soft_learnings):
    y = 2.4 + i * 1.52
    add_textbox(slide, 10.5, y, 8.8, 0.38, f'▶  {title}', 'Tahoma', 14, True, DARK_BLUE)
    add_textbox(slide, 10.5, y+0.38, 8.8, 0.75, desc, 'Calibri', 12, False, GRAY_TEXT)

add_footer(slide)

# ══════════════════════════════════════════════════════════════════
# SLIDE 7: Results / Demo
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 20, 11.25, WHITE)
add_rect(slide, 0, 0, 20, 1.35, ATT_BLUE)
add_rect(slide, 0, 0, 0.25, 1.35, ATT_YELLOW)
add_textbox(slide, 0.6, 0.3, 18, 0.75, '成果总结', 'Tahoma', 32, True, WHITE)

# Big stats row
stats2 = [
    ('5', '难度等级\nD1→D5'),
    ('2', '游戏模式\n闯关+节奏'),
    ('43KB', '单文件\n零依赖'),
    ('6', '可解锁技能\n元成长系统'),
    ('∞', '随机关卡\n程序化生成'),
]
for i, (num, lbl) in enumerate(stats2):
    add_stat_callout(slide, 0.4 + i*3.84, 1.6, 3.4, 1.9, num, lbl, ATT_BLUE)

# Deliverables
add_textbox(slide, 0.5, 3.85, 19, 0.5,
            '交付物', 'Tahoma', 20, True, DARK_BLUE)

deliverables = [
    ('🎮  platform_level/game.html',
     '完整横版2D闯关游戏 · 5难度 · 随机关卡 · 元成长 · 双攻击模式'),
    ('🎵  rhythm/rhythm.html',
     '节奏音游模块 · MIDI音频分析 · 独立可玩'),
    ('🏠  index.html',
     '游戏中心入口 · 整合两款游戏 · 最高分展示 · 粒子动效背景'),
    ('📁  platform_level/level_medium.json',
     '预设关卡配置文件 · 支持自定义关卡编辑'),
]
for i, (title, desc) in enumerate(deliverables):
    y = 4.55 + i * 1.45
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, 0.4, y, 19.2, 1.2, bg)
    add_rect(slide, 0.4, y, 0.12, 1.2, ATT_BLUE if i % 2 == 0 else ATT_YELLOW)
    add_textbox(slide, 0.7, y+0.08, 7, 0.45, title, 'Tahoma', 14, True, DARK_BLUE)
    add_textbox(slide, 0.7, y+0.55, 18, 0.5, desc, 'Calibri', 13, False, GRAY_TEXT)

add_footer(slide)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8: Thank You
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 20, 11.25, ATT_BLUE)
add_rect(slide, 0, 0, 0.25, 11.25, ATT_YELLOW)
add_rect(slide, 19.75, 0, 0.25, 11.25, ATT_YELLOW)

add_textbox(slide, 1, 3.2, 18, 1.8,
            'THANK YOU', 'Tahoma', 66, True, WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 1, 5.2, 18, 0.7,
            '感谢聆听 · 欢迎交流', 'Calibri', 26, False, ATT_LBLUE, align=PP_ALIGN.CENTER)

add_rect(slide, 6, 6.2, 8, 0.04, ATT_YELLOW)

add_textbox(slide, 1, 6.6, 18, 0.5,
            'Kenny  ·  ATT Group  ·  2026',
            'Calibri', 18, False, RGBColor(0xCC,0xDD,0xEE), align=PP_ALIGN.CENTER)
add_textbox(slide, 1, 7.2, 18, 0.4,
            '35 Ubi Crescent, ATT Building, Singapore, 408585',
            'Calibri', 14, False, RGBColor(0x88,0xAA,0xCC), align=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 10.65, 19, 0.4,
            FOOTER_TXT, 'Calibri', 8, False, RGBColor(0x88,0xAA,0xCC),
            align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────
out = r'D:\claude_ai_game\ATT_游戏项目汇报_Kenny.pptx'
prs.save(out)
print(f'Saved: {out}')
